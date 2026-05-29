from pptx import Presentation
import sys

def read_deck():
    prs = Presentation('DAT_Pitch_Deck.pptx')
    for i, slide in enumerate(prs.slides):
        print(f"\n--- Slide {i} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                print(f"[{shape.name}]: {shape.text}")

if __name__ == "__main__":
    read_deck()
