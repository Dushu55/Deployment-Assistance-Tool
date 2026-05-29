from pptx import Presentation

def update_deck():
    prs = Presentation('DAT_Pitch_Deck.pptx')
    
    # Slide 0
    slide0 = prs.slides[0]
    for shape in slide0.shapes:
        if "Subtitle" in shape.name or (hasattr(shape, "text") and "Unifying" in shape.text):
            shape.text = "The Autonomous, Polyglot DevSecOps Orchestrator"
            
    # Slide 1
    slide1 = prs.slides[1]
    for shape in slide1.shapes:
        if "Content" in shape.name:
            shape.text = (
                "The Business Problem:\n"
                "• Tool Sprawl: Engineering teams juggle 20+ disjointed tools.\n"
                "• Alert Fatigue & False Positives: Devs are blocked by deep dependency vulnerabilities in unreachable (dead) code.\n"
                "• Passive Reporting: Tools find bugs but offer zero engineering automation to actually fix them.\n"
                "• Polyglot Fragmentation: Nodes, Python, Java, C#, Go, and Rust require entirely disjointed security stacks."
            )

    # Slide 2
    slide2 = prs.slides[2]
    for shape in slide2.shapes:
        if "Content" in shape.name:
            shape.text = (
                "One CLI. 20+ Security Engines. Full Autonomous Remediation.\n"
                "• Polyglot Routing: Automatically detects Node, Python, Java, C#, Go, and Rust ecosystems and provisions native tools.\n"
                "• Normalized Intelligence: Translates 20+ tool outputs into a single deployment matrix.\n"
                "• Automated Fixing: Upgrades DAT from a 'reporter' to a 'fixer' using deterministic AST rewrites and LLM code refactoring."
            )

    # Slide 4
    slide4 = prs.slides[4]
    for shape in slide4.shapes:
        if "Title" in shape.name:
            shape.text = "AI-Driven Autonomous Remediation"
        elif "Content" in shape.name:
            shape.text = (
                "Not Just Security Reporting—Self-Healing Infrastructure.\n"
                "Capabilities Included:\n"
                "• Deterministic Auto-Fixers: AST-Grep safely rewrites insecure code patterns (e.g., eval to JSON.parse) inline.\n"
                "• Auto-Distroless Refactoring: Google Gemini LLMs instantly rewrite vulnerable Dockerfiles into secure, multi-stage Distroless builds.\n"
                "• Test-Driven Rollback Loop: Automatically executes native test suites (npm, pytest, go test). If an AI fix fails tests, DAT silently reverts the code via Git.\n"
                "• Self-Healing PR Agent: DAT branches, commits, pushes, and opens fully formatted GitHub Pull Requests automatically."
            )

    # Slide 7
    slide7 = prs.slides[7]
    for shape in slide7.shapes:
        if "Title" in shape.name:
            shape.text = "V2 Enhancements: Now in Production"
        elif "Content" in shape.name:
            shape.text = (
                "Major Upgrades Released:\n"
                "• Reachability Analysis Engine: Scans native Call-Graphs (Node, Python, Rust, Java, C#) to prove if a CVE is actually reachable. Downgrades false positives to unblock developers.\n"
                "• Ephemeral Environment Scans: Hooks into Vercel/AWS to spin up live PR previews, runs dynamic OWASP ZAP & k6 tests, and destroys the environment.\n"
                "• Native GitHub App: Zero-YAML, webhook-driven 1-click install for organizational CI/CD automation."
            )

    prs.save('DAT_Pitch_Deck_V2.pptx')
    print("Successfully updated presentation to DAT_Pitch_Deck_V2.pptx")

if __name__ == "__main__":
    update_deck()
